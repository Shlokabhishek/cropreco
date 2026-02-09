"""
HTML to PowerPoint Converter
Converts HTML documents to PowerPoint presentations
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from bs4 import BeautifulSoup
import re

def clean_text(text):
    """Clean and normalize text"""
    if not text:
        return ""
    # Remove extra whitespace
    text = re.sub(r'\s+', ' ', text)
    return text.strip()

def hex_to_rgb(hex_color):
    """Convert hex color to RGB"""
    hex_color = hex_color.lstrip('#')
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

def add_title_slide(prs, title, subtitle=""):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title slide layout
    
    title_shape = slide.shapes.title
    title_shape.text = title
    
    if subtitle and len(slide.placeholders) > 1:
        subtitle_shape = slide.placeholders[1]
        subtitle_shape.text = subtitle
    
    return slide

def add_content_slide(prs, title, content_items, layout_type='bullet'):
    """Add a content slide with title and bullet points or content"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and content layout
    
    title_shape = slide.shapes.title
    title_shape.text = title
    
    # Add content
    if len(slide.placeholders) > 1:
        content_shape = slide.placeholders[1]
        text_frame = content_shape.text_frame
        text_frame.clear()
        
        for i, item in enumerate(content_items):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            p.text = item
            p.level = 0
            p.font.size = Pt(18)
    
    return slide

def add_two_column_slide(prs, title, left_items, right_items):
    """Add a slide with two columns"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title_shape = slide.shapes.title
    title_shape.text = title
    
    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5))
    text_frame = left_box.text_frame
    for i, item in enumerate(left_items):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(14)
    
    # Right column
    right_box = slide.shapes.add_textbox(Inches(5.5), Inches(1.5), Inches(4.5), Inches(5))
    text_frame = right_box.text_frame
    for i, item in enumerate(right_items):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(14)
    
    return slide

def add_table_slide(prs, title, headers, rows):
    """Add a slide with a table"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(28)
    title_frame.paragraphs[0].font.bold = True
    
    # Add table
    rows_count = len(rows) + 1  # +1 for header
    cols_count = len(headers)
    
    table = slide.shapes.add_table(rows_count, cols_count, Inches(0.5), Inches(1.5), Inches(9), Inches(4.5)).table
    
    # Set column widths
    for col_idx in range(cols_count):
        table.columns[col_idx].width = Inches(9 / cols_count)
    
    # Add headers
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(52, 73, 94)
        
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.bold = True
        paragraph.font.size = Pt(12)
        paragraph.font.color.rgb = RGBColor(255, 255, 255)
    
    # Add data rows
    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_data in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_data)
            
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.size = Pt(10)
    
    return slide

def parse_html_to_pptx(html_file, output_file):
    """Convert HTML file to PowerPoint presentation"""
    
    print(f"Converting {html_file} to {output_file}...")
    
    # Read HTML file
    with open(html_file, 'r', encoding='utf-8') as f:
        html_content = f.read()
    
    # Parse HTML
    soup = BeautifulSoup(html_content, 'html.parser')
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Get title
    title_tag = soup.find('title')
    main_title = title_tag.text if title_tag else "Presentation"
    
    # Get subtitle from meta info or first p in header
    subtitle = ""
    meta_info = soup.find(class_=['meta-info', 'author-info', 'document-header'])
    if meta_info:
        subtitle = clean_text(meta_info.get_text())[:200]
    
    # Add title slide
    add_title_slide(prs, main_title, subtitle)
    
    # Process content
    content_sections = []
    current_section = None
    
    # Find all h1, h2, h3 headers and their content
    for element in soup.find_all(['h1', 'h2', 'h3', 'div']):
        
        if element.name in ['h1', 'h2']:
            # Save previous section
            if current_section and current_section['content']:
                content_sections.append(current_section)
            
            # Start new section
            current_section = {
                'title': clean_text(element.get_text()),
                'content': [],
                'level': element.name
            }
        
        elif element.name == 'h3' and current_section:
            # Add as sub-item
            current_section['content'].append({
                'type': 'subtitle',
                'text': clean_text(element.get_text())
            })
        
        elif element.name == 'div' and current_section:
            # Check for special divs (paper-entry, section-box, etc.)
            if 'paper-summary' in element.get('class', []) or 'paper-entry' in element.get('class', []):
                # This is a paper summary
                paper_title = element.find(class_=['paper-title'])
                if paper_title:
                    current_section['content'].append({
                        'type': 'text',
                        'text': clean_text(paper_title.get_text())
                    })
                
                # Get key points from various boxes
                for box in element.find_all(class_=['section-box', 'key-findings', 'results-box']):
                    box_title = box.find(class_='section-title')
                    if box_title:
                        current_section['content'].append({
                            'type': 'subtitle',
                            'text': clean_text(box_title.get_text())
                        })
                    
                    # Get bullet points
                    for li in box.find_all('li'):
                        current_section['content'].append({
                            'type': 'bullet',
                            'text': clean_text(li.get_text())
                        })
    
    # Add last section
    if current_section and current_section['content']:
        content_sections.append(current_section)
    
    # If no sections found, try a simpler approach
    if len(content_sections) < 3:
        content_sections = []
        current_section = None
        
        for element in soup.find_all(['h1', 'h2', 'h3', 'p', 'ul', 'ol', 'table']):
            if element.name in ['h1', 'h2']:
                # Save previous section
                if current_section and current_section['content']:
                    content_sections.append(current_section)
                
                # Start new section
                current_section = {
                    'title': clean_text(element.get_text()),
                    'content': []
                }
            
            elif current_section:
                if element.name == 'h3':
                    current_section['content'].append(clean_text(element.get_text()))
                
                elif element.name == 'p':
                    text = clean_text(element.get_text())
                    if text and len(text) > 10:  # Skip very short texts
                        current_section['content'].append(text)
                
                elif element.name in ['ul', 'ol']:
                    for li in element.find_all('li', recursive=False):
                        text = clean_text(li.get_text())
                        if text:
                            current_section['content'].append(f"• {text}")
                
                elif element.name == 'table':
                    # Extract table data
                    headers = [clean_text(th.get_text()) for th in element.find_all('th')]
                    rows = []
                    for tr in element.find_all('tr')[1:]:  # Skip header row
                        row = [clean_text(td.get_text()) for td in tr.find_all('td')]
                        if row:
                            rows.append(row)
                    
                    if headers and rows:
                        current_section['content'].append({
                            'type': 'table',
                            'headers': headers,
                            'rows': rows
                        })
        
        # Add last section
        if current_section and current_section['content']:
            content_sections.append(current_section)
    
    # Create slides from sections
    for section in content_sections[:30]:  # Limit to 30 slides
        title = section['title']
        
        # Split content into chunks if too long
        content_items = []
        table_data = None
        
        for item in section['content']:
            if isinstance(item, dict):
                if item['type'] == 'table':
                    table_data = item
                else:
                    text = item.get('text', '')
                    if len(text) <= 150:
                        content_items.append(text)
                    else:
                        # Split long text
                        words = text.split()
                        chunk = []
                        for word in words:
                            chunk.append(word)
                            if len(' '.join(chunk)) > 120:
                                content_items.append(' '.join(chunk))
                                chunk = []
                        if chunk:
                            content_items.append(' '.join(chunk))
            else:
                if len(item) <= 150:
                    content_items.append(item)
                else:
                    # Truncate long items
                    content_items.append(item[:147] + "...")
        
        # Add slides
        if table_data:
            # Create table slide
            add_table_slide(prs, title, table_data['headers'], table_data['rows'][:10])
        elif len(content_items) > 0:
            # Create content slides (max 7 items per slide)
            for i in range(0, len(content_items), 7):
                slide_title = title if i == 0 else f"{title} (cont.)"
                add_content_slide(prs, slide_title, content_items[i:i+7])
    
    # Add final slide
    add_title_slide(prs, "Thank You", "Questions & Discussion")
    
    # Save presentation
    prs.save(output_file)
    print(f"✓ Successfully created {output_file}")
    print(f"  Total slides: {len(prs.slides)}")

def main():
    """Main conversion function"""
    
    # List of files to convert
    files_to_convert = [
        ('RESEARCH_PAPER.html', 'RESEARCH_PAPER.pptx'),
        ('LITERATURE_SURVEY.html', 'LITERATURE_SURVEY.pptx'),
        ('RESEARCH_PAPERS_SUMMARY.html', 'RESEARCH_PAPERS_SUMMARY.pptx'),
        ('CROP_RECOMMENDER_PRESENTATION.html', 'CROP_RECOMMENDER_PRESENTATION.pptx'),
    ]
    
    print("=" * 60)
    print("HTML to PowerPoint Converter")
    print("=" * 60)
    
    for html_file, pptx_file in files_to_convert:
        try:
            parse_html_to_pptx(html_file, pptx_file)
            print()
        except FileNotFoundError:
            print(f"✗ File not found: {html_file}")
            print()
        except Exception as e:
            print(f"✗ Error converting {html_file}: {str(e)}")
            print()
    
    print("=" * 60)
    print("Conversion complete!")
    print("=" * 60)

if __name__ == "__main__":
    main()
